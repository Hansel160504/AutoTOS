# AutoTOS AI module — Ollama backend
#
# ── Patch history ──────────────────────────────────────────────────────────
# v4.9  – Instruction trim, Bloom pattern relaxation.
# v4.10 – Qwen3-4B token + negation fixes.
# v4.11 – Term-def expansion, "which statement best" hard-reject,
#          TF context-framing rejection, instruction rewrite, open verb diversity.
# v4.12 – MCQ diversity + fallback validation pass.
# v4.13 – Pronoun starters allowed, Analyzing bloom relaxation.
#
# v4.14  (Training-inference alignment)
#       Fix-58 – AUTOTOS_INSTRUCTION exact copy of train_v15.py.
#       Fix-59 – "Which statement best" exempt for Understanding bloom.
#       Fix-60 – "How would you implement/justify" re-allowed.
#       Fix-61 – Open-ended answers expanded to 2 sentences.
#       Fix-62 – TF note softened to match training.
#
# v4.15  (Definition-pattern coverage)
#       Fix-63 – _TERM_DEF_QUESTION_RE: "Define X" and "Explain the meaning".
#
# v4.16  (Diversity + dataset cleanup)
#       Fix-64 – Same-first-word exempt: question-word starters (how/what/etc).
#       Fix-65 – Analyzing Bloom pattern: "how would you compare/examine/…".
#       Fix-66 – Purpose-question cap: "what_is_purpose" capped opener.
#       Fix-67 – Junk distractor guard (_has_junk_distractors).
#       Fix-68 – MCQ subtopic SW: add purpose/goal/aim/role/function stopwords.
#
# v4.17  (Quality fixes from test-output audit)
#       Fix-68 – _MCQ_SUBTOPIC_SW actually applied (was in notes, not code).
#       Fix-69 – MCQ meta-phrase ban ("in the context of this scenario").
#       Fix-70 – "behaves under load" lazy template ban.
#       Fix-71 – "How would you describe/explain" capped opener.
#       Fix-72 – "What does X involve?" added to term-def rejection.
#       Fix-73 – Junk answer_text template detection.
#       Fix-74 – Circular choice detection (choice mirrors question stem).
#       Fix-75 – _mcq_opener_hint indentation fixed.
#
# v4.18  (Item-count guarantee — BUG FIXES)
#       Fix-76 – ThreadPoolExecutor future.result() wrapped in try/except.
#       Fix-77 – Safety net fills any remaining None slots with placeholders.
#       Fix-78 – /generate and /generate_from_records endpoints emit WARNING
#                when len(records) < max_items.
#
# v4.19  (Fixed-mode short-circuit)
#       Fix-79 – generate_questions_from_records now properly short-circuits
#                to the fixed question bank when FIXED_MODE = True, bypassing
#                all AI generation logic entirely.
#       Fix-80 – get_fixed_question now converts letter answers (A/B/C/D) to
#                the matching choice text so the frontend renders correctly.
# ============================================================

import os
import re
import json
import fitz
import base64
import hashlib
import logging
import random
import tempfile
import threading
import time
from collections import OrderedDict, deque
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import lru_cache
from io import BytesIO
from typing import Any, Dict, List, Optional, Set, Tuple
from pathlib import Path
import requests
from docx import Document
from pptx import Presentation

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from starlette.concurrency import run_in_threadpool
from fastapi.middleware.cors import CORSMiddleware

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# =====================================================
# CONFIGURATION
# =====================================================
OLLAMA_BASE_URL = os.environ.get("OLLAMA_BASE_URL", "http://ollama:11434")
OLLAMA_MODEL    = os.environ.get("OLLAMA_MODEL", "autotos")
OLLAMA_TIMEOUT  = int(os.environ.get("OLLAMA_TIMEOUT", "300"))

GENERATION_WORKERS = int(os.environ.get("GENERATION_WORKERS", "1"))

MAX_RETURN    = 50_000
CHUNK_SIZE    = 600
CHUNK_OVERLAP = 20

CACHE_MAX_FILES = 2_000

BASE_DIR        = os.path.dirname(__file__)
CACHE_DIR       = os.path.join(BASE_DIR, ".extracted_cache")
MODEL_CACHE_DIR = os.path.join(BASE_DIR, ".model_cache")
os.makedirs(CACHE_DIR,       exist_ok=True)
os.makedirs(MODEL_CACHE_DIR, exist_ok=True)

logger.info("Ollama config: base_url=%s model=%s timeout=%ds",
            OLLAMA_BASE_URL, OLLAMA_MODEL, OLLAMA_TIMEOUT)

SESSION = requests.Session()

# =====================================================
# IN-MEMORY LRU CACHE
# =====================================================
class _LRUMemCache:
    def __init__(self, maxsize: int = 512) -> None:
        self._cache: OrderedDict = OrderedDict()
        self._maxsize = maxsize
        self._lock = threading.Lock()

    def get(self, key: str) -> Optional[Any]:
        with self._lock:
            if key not in self._cache:
                return None
            self._cache.move_to_end(key)
            return self._cache[key]

    def set(self, key: str, value: Any) -> None:
        with self._lock:
            if key in self._cache:
                self._cache.move_to_end(key)
            self._cache[key] = value
            if len(self._cache) > self._maxsize:
                self._cache.popitem(last=False)

_mem_cache = _LRUMemCache(maxsize=512)

# =====================================================
# FIXED QUESTION BANK MODE  (Fix-79, Fix-80)
# =====================================================
# Set FIXED_MODE = True  → all generation uses the pre-built JSON bank.
# Set FIXED_MODE = False → live AI generation via Ollama.
# =====================================================
FIXED_MODE = True

_FIXED_BANK_ROOT   = Path(__file__).resolve().parent.parent
_FIXED_BANK_PATH_A = _FIXED_BANK_ROOT / "1.json"
_FIXED_BANK_PATH_B = Path(__file__).resolve().parent / "1.json"
FIXED_BANK_PATH    = _FIXED_BANK_PATH_A if _FIXED_BANK_PATH_A.exists() else _FIXED_BANK_PATH_B

_FIXED_BANK_CACHE: Optional[dict] = None

_LETTER_TO_INDEX = {"A": 0, "B": 1, "C": 2, "D": 3}


def load_fixed_bank() -> dict:
    """Load and cache the fixed question bank keyed by item_number (int)."""
    global _FIXED_BANK_CACHE
    if _FIXED_BANK_CACHE is None:
        if not FIXED_BANK_PATH.exists():
            raise FileNotFoundError(f"Fixed bank not found: {FIXED_BANK_PATH}")
        with FIXED_BANK_PATH.open("r", encoding="utf-8") as f:
            data = json.load(f)
        _FIXED_BANK_CACHE = {
            int(item["item_number"]): item
            for item in data
            if "item_number" in item
        }
        logger.info("Fixed bank loaded: %d items from %s", len(_FIXED_BANK_CACHE), FIXED_BANK_PATH)
    return _FIXED_BANK_CACHE


def get_fixed_question(item_number: int) -> Optional[dict]:
    """
    Return a normalised question dict from the fixed bank.

    Fix-80: letter answers (A/B/C/D) are resolved to the matching choice text
    so the frontend and DOCX builder receive a consistent format.
    """
    try:
        bank = load_fixed_bank()
    except FileNotFoundError as exc:
        logger.error("Fixed bank unavailable: %s", exc)
        return None

    q = bank.get(int(item_number))
    if not q:
        logger.warning("Fixed bank: item_number=%d not found", item_number)
        return None

    qtype   = (q.get("type") or "mcq").lower().strip()
    choices = [str(c) for c in (q.get("choices") or [])]

    # Resolve letter answer → choice text (Fix-80)
    raw_answer = str(q.get("answer") or "").strip()
    answer_upper = raw_answer.upper()
    if answer_upper in _LETTER_TO_INDEX and choices:
        idx = _LETTER_TO_INDEX[answer_upper]
        resolved_answer = choices[idx] if idx < len(choices) else raw_answer
    else:
        resolved_answer = raw_answer

    fixed_q: dict = {
        "type":     qtype if qtype in ("mcq", "truefalse", "open_ended") else "mcq",
        "concept":  q.get("topic", ""),
        "bloom":    q.get("bloom", "Remembering"),
        "question": q.get("question", ""),
        "choices":  choices,
        "answer":   resolved_answer,
    }

    if fixed_q["type"] != "open_ended":
        fixed_q["answer_text"] = q.get("answer_text", "")

    return fixed_q


# =====================================================
# FIXED-MODE BATCH GENERATOR  (Fix-79)
# =====================================================

def _generate_fixed_batch(records, max_items=None) -> List[dict]:
    """
    When FIXED_MODE is True, return questions directly from the fixed bank.
    Slot index i (0-based) maps to item_number i+1 in the bank.
    Items beyond the bank size get a clearly-labelled placeholder.
    """
    if isinstance(records, dict):
        records = records.get("records") or records.get("topics") or []
    records = records or []

    limit = min(len(records), max_items) if max_items is not None else len(records)
    out: List[dict] = []

    try:
        bank = load_fixed_bank()
        bank_size = max(bank.keys()) if bank else 0
    except Exception as exc:
        logger.error("Cannot load fixed bank: %s", exc)
        bank_size = 0

    for i in range(limit):
        item_number = i + 1
        q = get_fixed_question(item_number)
        if q:
            out.append(q)
        else:
            # Cycle through bank if exam is larger than 50 items
            if bank_size > 0:
                cycled_num = ((item_number - 1) % bank_size) + 1
                q_cycled = get_fixed_question(cycled_num)
                if q_cycled:
                    logger.info(
                        "Fixed bank: item %d out of range, cycling to item %d",
                        item_number, cycled_num
                    )
                    out.append(q_cycled)
                    continue

            # Hard fallback placeholder
            rec = records[i] if i < len(records) else {}
            inp = rec.get("input", {}) if isinstance(rec, dict) else {}
            topic = inp.get("concept") or inp.get("topic") or "General"
            out.append({
                "type":     "mcq",
                "concept":  topic,
                "bloom":    "Remembering",
                "question": f"[ITEM {item_number} NOT IN FIXED BANK] — {topic}",
                "choices":  ["(Not found)"] * 4,
                "answer":   "",
                "answer_text": "Item not found in fixed bank. Please add it or switch to AI mode.",
                "_generation_failed": True,
            })

    logger.info("Fixed-mode batch: returned %d/%d items", len(out), limit)
    return out


# =====================================================
# TYPE / BLOOM NORMALISERS
# =====================================================
TYPE_MAP = {
    "mcq": "mcq", "truefalse": "tf", "true_false": "tf", "tf": "tf",
    "open_ended": "open", "open-ended": "open", "openended": "open", "open": "open",
}
OUT_TYPE_NORMALIZE = {
    "mcq": "mcq", "truefalse": "truefalse", "tf": "truefalse",
    "true_false": "truefalse", "open_ended": "open_ended",
    "open": "open_ended", "open-ended": "open_ended",
}
BLOOM_MAP_SINGLE = {
    "knowledge": "Remembering", "understand": "Understanding", "apply": "Applying",
    "analyze": "Analyzing", "evaluate": "Evaluating", "create": "Creating",
    "remembering": "Remembering", "understanding": "Understanding",
    "applying": "Applying", "analyzing": "Analyzing", "evaluating": "Evaluating",
    "creating": "Creating",
}
BLOOM_CYCLE = {
    "remembering": ["Remembering", "Understanding"],
    "applying":    ["Applying",    "Analyzing"],
    "creating":    ["Evaluating",  "Creating"],
}


def normalize_bloom(bloom: str, slot_index: int = 0) -> str:
    stripped = (bloom or "").strip()
    _CANONICAL = {"Remembering", "Understanding", "Applying",
                  "Analyzing", "Evaluating", "Creating"}
    if stripped in _CANONICAL:
        return stripped
    key = stripped.lower()
    cycle = BLOOM_CYCLE.get(key)
    if cycle:
        return cycle[slot_index % 2]
    return BLOOM_MAP_SINGLE.get(key, stripped or "Remembering")


def normalize_type(qtype: str) -> str:
    return TYPE_MAP.get((qtype or "").strip().lower(), "mcq")


def normalize_out_type(raw_type: str) -> str:
    return OUT_TYPE_NORMALIZE.get((raw_type or "").strip().lower(), raw_type or "mcq")


# =====================================================
# OLLAMA CONNECTIVITY CHECK
# =====================================================
_ollama_ready = False


def _check_ollama() -> bool:
    global _ollama_ready
    if FIXED_MODE:
        return False   # No need to check when hardcoded
    try:
        r = SESSION.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=10)
        if r.ok:
            models = [m.get("name", "") for m in r.json().get("models", [])]
            available = any(
                OLLAMA_MODEL == m or OLLAMA_MODEL == m.split(":")[0]
                for m in models
            )
            if available:
                logger.info("Ollama ready. Model '%s' found.", OLLAMA_MODEL)
                _ollama_ready = True
            else:
                logger.warning(
                    "Ollama is running but model '%s' is NOT found.", OLLAMA_MODEL
                )
        return _ollama_ready
    except Exception as e:
        logger.warning("Ollama not reachable: %s", e)
        return False


if not FIXED_MODE:
    _check_ollama()

# =====================================================
# DISK CACHE
# =====================================================
def _sha256_bytes(b: bytes) -> str:
    h = hashlib.sha256()
    h.update(b)
    return h.hexdigest()


def _cache_path_for_hash(h: str) -> str:
    return os.path.join(CACHE_DIR, f"{h}.txt")


def _prompt_hash_key(prompt: str, max_tokens: int,
                     temperature: float, num_ctx: int) -> str:
    key = f"{OLLAMA_MODEL}|{max_tokens}|{temperature:.4f}|{num_ctx}|{prompt}"
    return hashlib.sha256(key.encode("utf-8")).hexdigest()


def _model_cache_path(key: str) -> str:
    return os.path.join(MODEL_CACHE_DIR, f"{key}.json")


def read_model_cache(key: str) -> Optional[Any]:
    p = _model_cache_path(key)
    if os.path.exists(p):
        try:
            with open(p, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            try:
                os.remove(p)
            except Exception:
                pass
    return None


_cache_write_count = 0
_CLEANUP_EVERY = 50


def write_model_cache(key: str, obj: Any) -> None:
    global _cache_write_count
    target = _model_cache_path(key)
    try:
        fd, tmp_path = tempfile.mkstemp(dir=MODEL_CACHE_DIR, suffix=".tmp")
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as f:
                json.dump(obj, f, ensure_ascii=False)
            os.replace(tmp_path, target)
        except Exception:
            try:
                os.remove(tmp_path)
            except Exception:
                pass
            raise
    except Exception as e:
        logger.debug("write_model_cache failed: %s", e)
        return
    _cache_write_count += 1
    if _cache_write_count % _CLEANUP_EVERY == 0:
        _cleanup_model_cache()


def _cleanup_model_cache(max_files: int = CACHE_MAX_FILES) -> None:
    try:
        all_files = [
            os.path.join(MODEL_CACHE_DIR, f)
            for f in os.listdir(MODEL_CACHE_DIR)
            if f.endswith(".json")
        ]
        if len(all_files) <= max_files:
            return
        all_files.sort(key=lambda p: os.path.getmtime(p))
        to_delete = all_files[:len(all_files) - max_files]
        for path in to_delete:
            try:
                os.remove(path)
            except Exception:
                pass
        logger.info("Cache cleanup: removed %d old entries (%d remain).",
                    len(to_delete), max_files)
    except Exception as e:
        logger.warning("Cache cleanup failed (non-fatal): %s", e)


_cleanup_model_cache()

CACHE_HITS   = 0
CACHE_MISSES = 0

# =====================================================
# GENERATION PROGRESS TRACKER
# =====================================================
_gen_progress: Dict[str, Any] = {"current": 0, "total": 0, "active": False}
_gen_progress_lock = threading.Lock()

# =====================================================
# UTILITIES
# =====================================================
def clean_text(txt) -> str:
    if txt is None:
        return ""
    if not isinstance(txt, str):
        try:
            txt = str(txt)
        except Exception:
            return ""
    return re.sub(r"\s+", " ", txt).strip()


def sanitize_prompt(prompt: str) -> str:
    return (prompt or "").strip()


_AWKWARD_PHRASING_RULES = [
    (re.compile(r'^Which of the following best describes which\b', re.IGNORECASE), "Which"),
    (re.compile(r'^Which of the following best describes what\b',  re.IGNORECASE), "What"),
    (re.compile(r'^Which best describes what\b',                   re.IGNORECASE), "What"),
    (re.compile(r'^Which best describes which\b',                  re.IGNORECASE), "Which"),
]

_FILL_IN_BLANK_RE = re.compile(r'_{3,}')


def is_fill_in_the_blank(question: str) -> bool:
    if not question:
        return False
    return bool(_FILL_IN_BLANK_RE.search(question))


def normalize_mcq_answer(answer_value: str, choices: list) -> str:
    answer_value = answer_value.strip()
    if len(answer_value) == 1 and answer_value.upper() in ["A", "B", "C", "D"]:
        return answer_value.upper()
    answer_lower = answer_value.lower()
    for i, choice in enumerate(choices):
        if choice and answer_lower in choice.lower():
            return chr(65 + i)
    return answer_value


def _clean_question_phrasing(text: str) -> str:
    if not text:
        return text
    for pattern, replacement in _AWKWARD_PHRASING_RULES:
        if pattern.match(text):
            text = pattern.sub(replacement + " ", text, count=1).strip()
            if text:
                text = text[0].upper() + text[1:]
            break
    return text


_TF_QUESTION_PREFIX_RE = re.compile(r"^(true\s+or\s+false\s*[:\-]\s*)", re.IGNORECASE)
_ARTIFACT_DIGIT_RE = re.compile(r"\b\d+\s+(?=[A-Z])")


def strip_question_prefix(text: str, is_open_ended: bool = False) -> str:
    text = (text or "").strip()
    if not text:
        return text
    text = _ARTIFACT_DIGIT_RE.sub("", text).strip()
    if is_open_ended:
        return _clean_question_phrasing(text)
    text = _TF_QUESTION_PREFIX_RE.sub("", text).strip()
    text = re.sub(r"^it is true that\s+", "", text, flags=re.IGNORECASE).strip()
    if text:
        text = text[0].upper() + text[1:]
    return _clean_question_phrasing(text)


_ANSWER_TEXT_MAX_CHARS = 220


def _truncate_answer_text(text: str) -> str:
    if not text:
        return ""
    t = re.sub(r"\s+", " ", text).strip()
    t = re.sub(r'^(answer\s*[:\-]\s*)', '', t, flags=re.IGNORECASE).strip()
    sentences = re.split(r'(?<=[.!?])\s+', t)
    first = sentences[0].strip() if sentences else t.strip()
    if first and first[-1] not in ".!?":
        first = first + "."
    if len(first) > _ANSWER_TEXT_MAX_CHARS:
        truncated = first[:_ANSWER_TEXT_MAX_CHARS]
        last_space = truncated.rfind(" ")
        if last_space > 0:
            truncated = truncated[:last_space]
        return truncated.rstrip('.,;:') + "..."
    return first


def _truncate_open_answer(text: str) -> str:
    """Return up to 2 complete sentences (Fix-61 v4.14)."""
    if not text:
        return ""
    t = re.sub(r"\s+", " ", text).strip()
    t = re.sub(r'^(answer\s*[:\-]\s*)', '', t, flags=re.IGNORECASE).strip()
    if not t:
        return ""
    _ABR_RE = re.compile(
        r'\b(e\.g|i\.e|etc|vs|Mr|Mrs|Dr|Prof|Sr|Jr|St|approx|fig|no)\.\s',
        re.IGNORECASE
    )
    masked = _ABR_RE.sub(lambda m: m.group(0).replace('.', '\x00'), t)
    parts = re.split(r'(?<=[.!?])\s+(?=[A-Z])|(?<=[.!?])$', masked)
    parts = [p.replace('\x00', '.').strip() for p in parts if p.strip()]
    selected = parts[:2]
    result = " ".join(selected)
    if result and result[-1] not in ".!?":
        result = result + "."
    return result


# =====================================================
# CHUNKING
# =====================================================
_chunk_cache: Dict[str, List[str]] = {}


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE,
               overlap: int = CHUNK_OVERLAP) -> List[str]:
    if not text:
        return []
    text = text.strip()
    chunks = []
    start = 0
    step = max(1, chunk_size - overlap)
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            boundary = max(text.rfind(". ", start, end),
                           text.rfind("! ", start, end),
                           text.rfind("? ", start, end))
            if boundary > start + chunk_size // 2:
                end = boundary + 1
            else:
                word_boundary = text.rfind(" ", start, end)
                if word_boundary > start:
                    end = word_boundary
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += step
    return chunks


def _find_best_chunk_idx(chunks: List[str], topic: str) -> int:
    if not chunks or not topic:
        return 0
    topic_lower = topic.lower()
    for i, chunk in enumerate(chunks):
        if topic_lower in chunk.lower():
            return i
    for word in topic_lower.split():
        if len(word) > 3:
            for i, chunk in enumerate(chunks):
                if word in chunk.lower():
                    return i
    return 0


def get_chunks_for_text(full_text: str) -> List[str]:
    key = hashlib.md5(full_text[:4096].encode("utf-8", errors="ignore")).hexdigest()
    if key not in _chunk_cache:
        _chunk_cache[key] = chunk_text(full_text)
        logger.info("Chunked document: %d chars -> %d chunks",
                    len(full_text), len(_chunk_cache[key]))
    return _chunk_cache[key]


# =====================================================
# FILE EXTRACTION
# =====================================================
def extract_text_from_bytes(file_bytes: bytes, filetype: str) -> str:
    text = ""
    try:
        if filetype == "pdf":
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            parts = []
            for page in doc:
                try:
                    ptext = page.get_text("text")
                    if ptext and ptext.strip():
                        parts.append(ptext)
                except Exception:
                    continue
            text = " ".join(parts)
        elif filetype == "docx":
            d = Document(BytesIO(file_bytes))
            text = " ".join(p.text for p in d.paragraphs if p.text and p.text.strip())
        elif filetype == "pptx":
            prs = Presentation(BytesIO(file_bytes))
            slt = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        slt.append(shape.text)
            text = " ".join(slt)
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning("Extraction error: %s", e)
    return clean_text(text)


def extract_text_from_path(path: str, max_chars: int = MAX_RETURN) -> str:
    if not path or not os.path.exists(path):
        return ""
    try:
        with open(path, "rb") as f:
            b = f.read()
    except Exception as e:
        logger.warning("Failed to read %s: %s", path, e)
        return ""
    h = _sha256_bytes(b)
    cache_file = _cache_path_for_hash(h)
    if os.path.exists(cache_file):
        try:
            with open(cache_file, "r", encoding="utf-8") as rf:
                return clean_text(rf.read())[:max_chars]
        except Exception:
            pass
    ext = (os.path.splitext(path)[1] or "").lower()
    filetype = {".pdf": "pdf", ".docx": "docx", ".doc": "docx",
                ".pptx": "pptx", ".ppt": "pptx"}.get(ext, "")
    extracted = extract_text_from_bytes(b, filetype)
    try:
        with open(cache_file, "w", encoding="utf-8") as wf:
            wf.write(extracted)
    except Exception:
        pass
    return clean_text(extracted)[:max_chars]


def lesson_from_upload(data_or_text: Optional[str]) -> str:
    if not data_or_text:
        return ""
    if isinstance(data_or_text, str) and os.path.exists(data_or_text):
        try:
            return extract_text_from_path(data_or_text, max_chars=MAX_RETURN)
        except Exception as e:
            logger.warning("Error extracting %s: %s", data_or_text, e)
            return ""
    if isinstance(data_or_text, str) and data_or_text.startswith("data:"):
        try:
            header, encoded = data_or_text.split(",", 1)
            file_bytes = base64.b64decode(encoded)
            ft = ("pdf"  if "pdf"  in header else
                  "docx" if ("docx" in header or "word" in header) else
                  "pptx" if ("pptx" in header or "presentation" in header) else "")
            return extract_text_from_bytes(file_bytes, ft)[:MAX_RETURN]
        except Exception as e:
            logger.warning("Base64 decode fail: %s", e)
            return ""
    try:
        return clean_text(data_or_text or "")[:MAX_RETURN]
    except Exception:
        return ""


# =====================================================
# PROMPT BUILDER  (used only when FIXED_MODE = False)
# =====================================================

AUTOTOS_INSTRUCTION = (
    "You are AutoTOS. Output ONLY valid JSON.\n\n"
    "CONSTRAINTS:\n"
    "- NO definition questions ('What is the term').\n"
    "- NO 'How does' or 'How will' (unless used for application or creation).\n"
    "- NO meta-phrases like 'In the context of...' or 'Assuming that...'.\n"
    "- Avoid misleading or tricky negative phrasing in True/False statements.\n"
    "- MCQ choices must begin with different words. DO NOT include A/B/C/D prefixes.\n\n"
    "BLOOM'S LEVEL STARTERS:\n"
    "- Remember/Understand: 'What is the primary purpose of...', 'Why is X important for...'\n"
    "- Apply/Analyze: 'How would you implement...', 'Given X, which approach...'\n"
    "- Evaluate/Create: 'How would you justify...', 'Which approach is most effective for...'\n\n"
    "MCQ JSON FORMAT:\n"
    "{\"type\":\"mcq\",\"concept\":\"...\",\"bloom\":\"...\",\"question\":\"...\","
    "\"choices\":[\"...\",\"...\",\"...\",\"...\"],\"answer\":\"A|B|C|D\","
    "\"answer_text\":\"Exactly 1 sentence why correct.\"}\n\n"
    "TF JSON FORMAT:\n"
    "{\"type\":\"truefalse\",\"concept\":\"...\",\"bloom\":\"...\",\"question\":\"...\","
    "\"answer\":\"true|false\",\"answer_text\":\"Exactly 1 sentence why correct.\"}"
)

AUTOTOS_INSTRUCTION_OPEN = (
    "You are AutoTOS. Output ONLY valid JSON.\n\n"
    "CONSTRAINTS:\n"
    "- NO definition questions.\n"
    "- The 'answer' string MUST be exactly 2 complete sentences ending in a period.\n\n"
    "OPEN-ENDED JSON FORMAT:\n"
    "{\"type\":\"open_ended\",\"concept\":\"...\",\"bloom\":\"...\",\"question\":\"...\","
    "\"answer\":\"<Exactly 2 complete sentences.>\"}"
)


def _build_avoid_block(seen_questions: List[str]) -> str:
    if not seen_questions:
        return ""
    recent = seen_questions[-3:]
    items  = "; ".join(f'"{q[:25]}"' for q in recent)
    return f"\n[Avoid repeating: {items}]\n"


def build_training_prompt(
    instruction: str,
    qtype: str,
    bloom: str,
    concept: str,
    context: str,
    extra_note: str = "",
    attempt_note: str = "",
    avoid_questions: Optional[List[str]] = None,
    attempt: int = 1,
) -> str:
    avoid_block = _build_avoid_block(avoid_questions or [])
    tf_note = ""
    if qtype == "tf":
        tf_note = (
            "\nTF RULES: Write a POSITIVE declarative statement (no question mark). "
            "Avoid misleading or tricky negative phrasing.\n"
        )
    retry_line = f"\n{attempt_note.strip()}\n" if attempt_note and attempt_note.strip() else ""
    ctx_suffix = f"{tf_note}{retry_line}{avoid_block}"
    user_msg = (
        "### Target Specification:\n"
        f"- Question Type: {qtype}\n"
        f"- Bloom's Level: {bloom}\n"
        f"- Concept: {concept}\n\n"
        "### Context (Source Material):\n"
        f"{context}{ctx_suffix}"
    )
    return (
        f"<|im_start|>system\n{instruction}<|im_end|>\n"
        f"<|im_start|>user\n{user_msg}<|im_end|>\n"
        f"<|im_start|>assistant\n"
    )


# =====================================================
# JSON EXTRACTORS
# =====================================================
def _extract_first_json(text: str) -> Optional[str]:
    if not text:
        return None
    start = text.find('{')
    if start == -1:
        return None
    depth = 0
    in_string = False
    escape = False
    for i, c in enumerate(text[start:], start):
        if escape:
            escape = False
            continue
        if c == '\\' and in_string:
            escape = True
            continue
        if c == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if c == '{':
            depth += 1
        elif c == '}':
            depth -= 1
            if depth == 0:
                return text[start:i + 1]
    return text[start:]


def _try_parse_json(json_str: str) -> Optional[Any]:
    if not json_str or not isinstance(json_str, str):
        return None
    try:
        return json.loads(json_str)
    except Exception:
        pass
    cleaned = re.sub(r',\s*([}\]])', r'\1', json_str)
    if cleaned.count('"') % 2 == 1:
        cleaned = cleaned + '"'
    open_braces = cleaned.count('{') - cleaned.count('}')
    if 0 < open_braces <= 6:
        try:
            return json.loads(cleaned + ('}' * open_braces))
        except Exception:
            pass
    try:
        return json.loads(cleaned)
    except Exception:
        return None


# =====================================================
# MODEL CALLER  (used only when FIXED_MODE = False)
# =====================================================
_NUM_CTX: Dict[str, int] = {"mcq": 1024, "tf": 1024, "open": 1024}
MAX_TOKENS_SINGLE: Dict[str, int] = {"mcq": 160, "tf": 85, "open": 120}


def ask_model(prompt: str, max_tokens: int = 200,
              temperature: float = 0.45,
              num_ctx: int = 1024) -> Optional[dict]:
    global CACHE_HITS, CACHE_MISSES
    if FIXED_MODE:
        logger.warning("ask_model called while FIXED_MODE=True — returning None.")
        return None
    if not _ollama_ready and not _check_ollama():
        logger.error("Ollama not ready.")
        return None

    prompt    = sanitize_prompt(prompt)
    cache_key = _prompt_hash_key(prompt, max_tokens, temperature, num_ctx)

    mem_hit = _mem_cache.get(cache_key)
    if mem_hit is not None:
        CACHE_HITS += 1
        return mem_hit if isinstance(mem_hit, dict) else None

    cached = read_model_cache(cache_key)
    if cached is not None:
        CACHE_HITS += 1
        _mem_cache.set(cache_key, cached)
        return cached if isinstance(cached, dict) else None
    CACHE_MISSES += 1

    payload = {
        "model":  OLLAMA_MODEL,
        "prompt": prompt,
        "stream": False,
        "format": "json",
        "options": {
            "temperature": temperature,
            "num_predict": max_tokens,
            "num_ctx":     num_ctx,
            "top_p":       0.95,
            "think":       False,
            "stop":        ["<|im_start|>", "<|im_end|>", "<|endoftext|>"],
        },
    }

    try:
        start = time.time()
        resp  = SESSION.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json=payload,
            timeout=OLLAMA_TIMEOUT,
        )
        duration = time.time() - start
        if not resp.ok:
            logger.warning("Ollama returned %d: %s", resp.status_code, resp.text[:200])
            return None

        raw = resp.json().get("response", "")
        logger.info("ask_model duration=%.2fs raw_len=%d", duration, len(raw))
        if not raw or not raw.strip():
            return None

        raw = re.sub(r"<think>[\s\S]*?</think>", "", raw).strip()
        if not raw:
            return None

        json_str = _extract_first_json(raw)
        parsed   = _try_parse_json(json_str) if json_str else None
        if parsed is None:
            logger.warning("No JSON found in: %s", raw[:200])
            return None

        write_model_cache(cache_key, parsed)
        _mem_cache.set(cache_key, parsed)
        return parsed

    except requests.exceptions.Timeout:
        logger.error("Ollama request timed out after %ds", OLLAMA_TIMEOUT)
        return None
    except Exception as e:
        logger.exception("ask_model error: %s", e)
        return None


def _warmup_model() -> None:
    if FIXED_MODE:
        return
    try:
        resp = SESSION.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json={
                "model":  OLLAMA_MODEL,
                "prompt": "Say OK",
                "stream": False,
                "options": {"num_predict": 3, "temperature": 0.0, "num_ctx": 1024},
            },
            timeout=OLLAMA_TIMEOUT,
        )
        if resp.ok:
            logger.info("Warm-up done — model loaded into memory.")
    except Exception as e:
        logger.warning("Warm-up ping failed (non-fatal): %s", e)


if not FIXED_MODE and _ollama_ready:
    logger.info("Warming up Ollama (model load ping)...")
    _warmup_model()


# =====================================================
# NORMALISATION / KEY MAPPING
# =====================================================
def _normalize_output_keys(out: Dict[str, Any]) -> Dict[str, Any]:
    if not isinstance(out, dict):
        return {}
    mapping = {
        "statement": "question", "prompt": "question",
        "sample_answer": "answer_text", "sample_response": "answer_text",
        "model_answer": "answer", "explanation": "answer_text",
        "solution": "answer_text", "rationale": "answer_text",
        "ans": "answer", "correct": "answer",
    }
    for src, dst in mapping.items():
        if src in out and dst not in out:
            out[dst] = out[src]
    if "answer" not in out and "sample_answer" in out:
        out["answer"] = out["sample_answer"]
    if isinstance(out.get("answer"), bool):
        out["answer"] = "true" if out["answer"] else "false"
    return out


_CHOICE_LETTER_PREFIX_RE = re.compile(
    r'^(?:\([A-Da-d]\)|[A-Da-d][).:\-]\s*|[A-Da-d]\s+(?=[A-Z]))'
)


def _strip_choice_letter_prefix(text: str) -> str:
    if not text:
        return text
    stripped = _CHOICE_LETTER_PREFIX_RE.sub("", text).strip()
    if stripped == text and len(text) > 2 and text[0].upper() in "ABCD" and text[1] == " ":
        stripped = text[2:].strip()
    return stripped


def _choice_has_letter_prefix(text: str) -> bool:
    if not text or len(text) < 2:
        return False
    return bool(_CHOICE_LETTER_PREFIX_RE.match(text))


def normalize_generated_question(q: dict, expected_display_type: str,
                                  topic: str, bloom_level: str) -> dict:
    q = q or {}
    if not isinstance(q, dict):
        q = {"question": str(q)}
    q = _normalize_output_keys(q)
    raw_out_type = q.get("type") or expected_display_type
    display_type = normalize_out_type(raw_out_type) or expected_display_type
    out = {
        "type":    display_type,
        "concept": topic,
        "bloom":   bloom_level,
        "question": strip_question_prefix(
            clean_text(q.get("question") or q.get("prompt") or ""),
            is_open_ended=(display_type == "open_ended")
        ),
        "choices":  [],
        "answer":   "",
    }
    if display_type != "open_ended":
        raw_answer_text = clean_text(q.get("answer_text") or q.get("explanation") or "")
        out["answer_text"] = _truncate_answer_text(raw_answer_text)

    choices_raw = q.get("choices") or q.get("options")
    if isinstance(choices_raw, dict):
        keys = sorted(choices_raw.keys(), key=lambda s: s.upper())
        out["choices"] = [
            _strip_choice_letter_prefix(clean_text(choices_raw[k])) for k in keys
        ][:4]
    elif isinstance(choices_raw, list):
        out["choices"] = [
            _strip_choice_letter_prefix(clean_text(x)) for x in choices_raw
        ][:4]

    ans = q.get("answer", "")
    if isinstance(ans, bool):
        out["answer"] = "true" if ans else "false"
    elif isinstance(ans, (int, float)):
        idx = int(ans)
        out["answer"] = (out["choices"][idx] if out["choices"] and 0 <= idx < len(out["choices"]) else str(ans))
    elif isinstance(ans, str):
        a = ans.strip()
        a_stripped = _CHOICE_LETTER_PREFIX_RE.sub("", a).strip()
        if display_type == "mcq":
            if re.fullmatch(r"[A-Da-d]", a.strip()) and out["choices"]:
                idx = ord(a.strip().upper()) - ord("A")
                out["answer"] = (out["choices"][idx] if 0 <= idx < len(out["choices"]) else a.strip())
            elif re.match(r"^[A-Da-d][).:\s]", a):
                matched = next((c for c in out["choices"] if c.lower() == a_stripped.lower()), None)
                out["answer"] = matched or a_stripped
            else:
                matched = next((c for c in out["choices"] if c.lower() == a.lower()), None)
                out["answer"] = matched or a
        elif display_type == "truefalse":
            a_lower = a.lower().rstrip(".")
            if a_lower in ("true", "false"):   out["answer"] = a_lower
            elif a_lower in ("1", "yes"):       out["answer"] = "true"
            elif a_lower in ("0", "no"):        out["answer"] = "false"
            else:                               out["answer"] = a
        elif display_type == "open_ended":
            out["answer"] = _truncate_open_answer(clean_text(a))
        else:
            out["answer"] = a
    else:
        raw_ans = clean_text(str(ans or ""))
        if display_type == "open_ended":
            out["answer"] = _truncate_open_answer(raw_ans)
        else:
            out["answer"] = raw_ans

    if display_type == "open_ended" and out["answer"]:
        out["answer"] = _truncate_open_answer(out["answer"])

    return out


# =====================================================
# FINGERPRINTING  (used in AI mode only)
# =====================================================
_FP_STOPWORDS = re.compile(
    r'\b(a|an|the|and|or|of|in|on|to|for|is|are|was|were|by|at|be|its|it|'
    r'that|this|these|those|with|as|from|into|about|which|how|what|does|do)\b')


def _question_fingerprint(q: dict) -> str:
    raw = (q.get("question") or "").lower().strip()
    raw = re.sub(r"[^\w\s]", " ", raw)
    raw = _FP_STOPWORDS.sub(" ", raw)
    raw = re.sub(r"\s+", " ", raw).strip()
    words = raw.split()
    words = [w[:-1] if w.endswith("s") and len(w) > 4 else w for w in words]
    raw   = " ".join(words)
    qtext   = raw[:35]
    concept = re.sub(r"\s+", "_", (q.get("concept") or "").lower().strip())
    return f"{concept}::{qtext}"


def _question_stem(q: dict) -> str:
    return (q.get("question") or "")[:80].strip()


# =====================================================
# VALIDATORS  (used in AI mode only)
# =====================================================
_OPEN_PLACEHOLDERS = {
    "model answer here", "model answer here.", "answer here",
    "write answer here", "<complete model answer based on context>",
}


def _is_valid_answer(q: dict, display_type: str) -> bool:
    answer  = clean_text(str(q.get("answer") or "")).strip()
    ans_low = answer.lower().rstrip(".")
    if ans_low in {"—", "-", "", "answer:"}:
        return False
    if display_type == "mcq":
        choices = q.get("choices") or []
        if not answer or len(choices) != 4:
            return False
        for c in choices:
            if not c or not c.strip():
                return False
        choices_lower = [c.lower().strip() for c in choices]
        answer_lower  = answer.lower().strip()
        if answer_lower not in choices_lower:
            matched = any(
                answer_lower in c or c in answer_lower
                for c in choices_lower
                if len(c) > 5 and len(answer_lower) > 5
            )
            if not matched:
                return False
        if len(choices_lower) != len(set(choices_lower)):
            return False
        return True
    elif display_type == "truefalse":
        return ans_low in ("true", "false")
    elif display_type == "open_ended":
        if len(answer) < 15:
            return False
        if ans_low in _OPEN_PLACEHOLDERS:
            return False
        return True
    return True


def _is_valid_fallback_question(q: dict, display_type: str) -> bool:
    qtext = (q.get("question") or "").strip()
    if not qtext or len(qtext) < 20:
        return False
    if display_type == "mcq":
        choices = q.get("choices") or []
        if len(choices) != 4:
            return False
        if not (q.get("answer") or "").strip():
            return False
    elif display_type == "truefalse":
        answer = (q.get("answer") or "").strip().lower().rstrip(".")
        if answer not in ("true", "false"):
            return False
    return True


# =====================================================
# SINGLE AI QUESTION GENERATOR  (FIXED_MODE = False only)
# =====================================================
_RETRY_NOTES = [
    "",
    "Try a different angle — focus on a specific component or detail.",
    "Use a concrete scenario or example from the context.",
]
_MAX_ATTEMPTS = 3


def _generate_single(
    topic, prompt_type, display_type, bloom, context, max_tok,
    seen_fps, record_idx,
    seen_questions=None, fp_lock=None,
    seen_answer_fps=None, answer_fp_lock=None,
    **kwargs,
):
    """Generate a single question via Ollama (AI mode only)."""
    instruction = AUTOTOS_INSTRUCTION_OPEN if prompt_type == "open" else AUTOTOS_INSTRUCTION
    ctx_size    = _NUM_CTX.get(prompt_type, 1024)
    avoid_list  = list((seen_questions or []))[-3:]

    for attempt in range(1, _MAX_ATTEMPTS + 1):
        temp         = min(0.80, 0.45 + 0.175 * (attempt - 1))
        attempt_note = _RETRY_NOTES[min(attempt - 1, len(_RETRY_NOTES) - 1)]

        prompt = build_training_prompt(
            instruction, prompt_type, bloom, topic, context,
            attempt_note=attempt_note,
            avoid_questions=avoid_list,
            attempt=attempt,
        )

        generated = ask_model(prompt, max_tokens=max_tok, temperature=temp, num_ctx=ctx_size)
        if generated is None:
            time.sleep(0.05 * attempt)
            continue

        generated   = _normalize_output_keys(generated)
        candidate_q = normalize_generated_question(generated, display_type, topic, bloom)
        qtext       = (candidate_q.get("question") or "").strip()
        if not qtext:
            time.sleep(0.05 * attempt)
            continue

        fp     = _question_fingerprint(candidate_q)
        is_dup = False
        if fp_lock:
            with fp_lock:
                if fp in seen_fps:
                    is_dup = True
                else:
                    seen_fps.add(fp)
        else:
            if fp in seen_fps:
                is_dup = True
            else:
                seen_fps.add(fp)
        if is_dup:
            avoid_list.append(qtext[:25])
            avoid_list = avoid_list[-3:]
            time.sleep(0.05 * attempt)
            continue

        if not _is_valid_answer(candidate_q, display_type):
            avoid_list.append(qtext[:25])
            avoid_list = avoid_list[-3:]
            time.sleep(0.05 * attempt)
            continue

        return candidate_q

    return None


# =====================================================
# BATCH GENERATOR  (Fix-79: FIXED_MODE short-circuits here)
# =====================================================

def generate_questions_from_records(records_or_topics, max_items=None) -> List[dict]:
    """
    Main entry point for batch question generation.

    When FIXED_MODE = True  → returns questions directly from the fixed JSON bank.
    When FIXED_MODE = False → runs the full Ollama AI generation pipeline.
    """
    try:
        if isinstance(records_or_topics, dict):
            records = records_or_topics.get("records") or records_or_topics.get("topics") or []
        else:
            records = records_or_topics or []

        if not isinstance(records, list):
            return []

        # ── Fix-79: hard short-circuit for fixed mode ──────────────
        if FIXED_MODE:
            logger.info(
                "FIXED_MODE=True: serving %d items from fixed question bank.",
                min(len(records), max_items) if max_items is not None else len(records)
            )
            return _generate_fixed_batch(records, max_items=max_items)

        # ── AI generation pipeline (FIXED_MODE = False) ─────────────
        limit = min(len(records), max_items) if max_items is not None else len(records)

        topic_slot_counter: Dict[str, List[int]] = {}
        slots = []

        for i in range(limit):
            rec = records[i]
            input_obj = rec.get("input", {}) if isinstance(rec, dict) else {}
            topic = (
                input_obj.get("concept")
                or input_obj.get("topic")
                or rec.get("instruction", "General")
                or "General"
            )
            raw_bloom = (
                input_obj.get("bloom")
                or (rec.get("output", {}) or {}).get("bloom")
                or "Remembering"
            )
            raw_type = (
                input_obj.get("type")
                or (rec.get("output", {}) or {}).get("type")
                or "mcq"
            )
            prompt_type  = normalize_type(raw_type)
            display_type = normalize_out_type(
                {"mcq": "mcq", "tf": "truefalse", "open": "open_ended"}.get(
                    prompt_type, prompt_type
                )
            )
            bloom = normalize_bloom(raw_bloom, slot_index=i)
            candidate = (
                input_obj.get("context")
                or input_obj.get("learn_material")
                or input_obj.get("file_path")
                or rec.get("file_path")
                or ""
            )
            full_text = lesson_from_upload(candidate) if candidate else ""
            context   = ""

            if full_text:
                chunks    = get_chunks_for_text(full_text)
                text_hash = hashlib.md5(
                    full_text[:4096].encode("utf-8", errors="ignore")
                ).hexdigest()[:8]
                topic_key = f"{topic}::{text_hash}"
                base_idx  = _find_best_chunk_idx(chunks, topic)

                if topic_key not in topic_slot_counter:
                    idxs = list(range(len(chunks)))
                    if base_idx in idxs:
                        idxs.remove(base_idx)
                    idxs.insert(0, base_idx)
                    tail = idxs[1:]
                    random.shuffle(tail)
                    idxs[1:] = tail
                    topic_slot_counter[topic_key] = idxs

                idxs = topic_slot_counter[topic_key]
                chunk_idx = idxs.pop(0)
                if not idxs:
                    topic_slot_counter.pop(topic_key, None)
                context = chunks[chunk_idx]
            else:
                logger.warning("record=%d topic=%r has NO learning material.", i + 1, topic)

            slots.append({
                "record_idx":  i,
                "topic":       topic,
                "bloom":       bloom,
                "prompt_type": prompt_type,
                "display_type": display_type,
                "context":     context,
                "record":      rec,
            })

        _fp_lock        = threading.Lock()
        _stems_lock     = threading.Lock()
        seen_fps:   set = set()
        seen_stems: deque = deque(maxlen=16)

        with _gen_progress_lock:
            _gen_progress["current"] = 0
            _gen_progress["total"]   = len(slots)
            _gen_progress["active"]  = True

        def _run_slot(slot: dict):
            with _stems_lock:
                stems_snapshot = list(seen_stems)
            q = _generate_single(
                slot["topic"], slot["prompt_type"], slot["display_type"],
                slot["bloom"], slot["context"],
                MAX_TOKENS_SINGLE.get(slot["prompt_type"], 115),
                seen_fps, slot["record_idx"] + 1,
                seen_questions=stems_snapshot,
                fp_lock=_fp_lock,
            )
            with _gen_progress_lock:
                _gen_progress["current"] += 1
            if q is not None:
                stem = _question_stem(q)
                if stem:
                    with _stems_lock:
                        seen_stems.append(stem)
            return slot["record_idx"], slot, q

        workers = min(GENERATION_WORKERS, len(slots))
        logger.info("AI mode: generating %d questions with %d worker(s)", len(slots), workers)
        slot_start = time.time()
        results: List[Optional[tuple]] = [None] * len(slots)

        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_map = {executor.submit(_run_slot, slot): i for i, slot in enumerate(slots)}
            for future in as_completed(future_map):
                try:
                    orig_i, slot, q = future.result()
                    results[orig_i] = (slot, q)
                except Exception as exc:
                    orig_i = future_map[future]
                    logger.error(
                        "Fix-76: worker crashed for slot %d — error: %s", orig_i, exc
                    )

        # Fix-77: safety net
        for i, res in enumerate(results):
            if res is None:
                logger.warning("Fix-77: slot %d has no result — filling with placeholder.", i)
                results[i] = (slots[i], None)

        elapsed = time.time() - slot_start
        logger.info("All %d questions generated in %.1fs", len(slots), elapsed)

        with _gen_progress_lock:
            _gen_progress["active"] = False

        out_questions = []
        for slot, q in results:
            if q is not None:
                out_questions.append(q)
            else:
                rec      = slot["record"]
                fallback = rec.get("output") if isinstance(rec, dict) else None
                used_fallback = False
                if fallback and isinstance(fallback, dict):
                    normalized_fb = normalize_generated_question(
                        fallback, slot["display_type"], slot["topic"], slot["bloom"]
                    )
                    if _is_valid_fallback_question(normalized_fb, slot["display_type"]):
                        out_questions.append(normalized_fb)
                        used_fallback = True
                if not used_fallback:
                    out_questions.append({
                        "type":    slot["display_type"],
                        "concept": slot["topic"],
                        "bloom":   slot["bloom"],
                        "question": f"[GENERATION FAILED] Review this item — {slot['topic']}",
                        "choices": (["(Generation failed)"] * 4 if slot["display_type"] == "mcq" else []),
                        "answer": "",
                        "answer_text": "Generation failed. Please delete or replace.",
                        "_generation_failed": True,
                    })
        return out_questions

    except Exception as e:
        logger.exception("generate_questions_from_records error: %s", e)
        return []


def generate_quiz_for_topics(records_or_topics, max_items=None, test_labels=None, *args, **kwargs):
    try:
        quizzes = generate_questions_from_records(records_or_topics, max_items)
    except Exception as e:
        logger.exception("generate_questions_from_records error: %s", e)
        quizzes = []

    if test_labels and isinstance(test_labels, (list, tuple)):
        for idx, item in enumerate(quizzes):
            if isinstance(item, dict):
                item["test_header"] = test_labels[idx] if idx < len(test_labels) else ""
    return {"quizzes": quizzes}


# =====================================================
# CACHE STATS
# =====================================================
def get_model_cache_stats():
    cache_files = sum(1 for f in os.listdir(MODEL_CACHE_DIR) if f.endswith(".json"))
    return {
        "fixed_mode":       FIXED_MODE,
        "cache_hits":       CACHE_HITS,
        "cache_misses":     CACHE_MISSES,
        "mem_cache_size":   len(_mem_cache._cache),
        "disk_cache_files": cache_files,
    }


def load_jsonl(path):
    records = []
    with open(path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                records.append(json.loads(line))
            except Exception as e:
                logger.warning("Skipping bad line: %s", e)
    return records


# =====================================================
# FASTAPI APP
# =====================================================
app = FastAPI(title="AutoTOS AI Service", version="4.19-fixed")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)


class GenerateRequest(BaseModel):
    records:     List[dict]
    max_items:   Optional[int]       = None
    test_labels: Optional[List[str]] = None


class ExtractRequest(BaseModel):
    data: str


@app.get("/health")
async def health():
    fixed = FIXED_MODE
    bank_ok = False
    bank_size = 0
    if fixed:
        try:
            bank = load_fixed_bank()
            bank_ok   = True
            bank_size = len(bank)
        except Exception:
            pass
    ready = fixed or _ollama_ready or _check_ollama()
    return {
        "status":            "ok" if ready else "degraded",
        "fixed_mode":        fixed,
        "fixed_bank_ok":     bank_ok,
        "fixed_bank_size":   bank_size,
        "ollama_ready":      _ollama_ready,
        "ollama_model":      OLLAMA_MODEL,
        "chunk_size":        CHUNK_SIZE,
        "generation_workers": GENERATION_WORKERS,
        "version":           "4.19",
    }


@app.get("/cache_stats")
async def cache_stats():
    return get_model_cache_stats()


@app.get("/progress")
async def generation_progress():
    with _gen_progress_lock:
        return dict(_gen_progress)


@app.post("/extract")
async def extract_text(req: ExtractRequest):
    try:
        text = await run_in_threadpool(lambda: lesson_from_upload(req.data))
        return {"text": text or ""}
    except Exception as e:
        logger.exception("extract error: %s", e)
        raise HTTPException(status_code=500, detail="Extraction failed")


@app.post("/generate")
async def generate(req: GenerateRequest):
    if not req.records:
        raise HTTPException(status_code=400, detail="records must be non-empty")
    if req.max_items is not None and len(req.records) < req.max_items:
        logger.warning(
            "Fix-78 /generate: got %d records but max_items=%d — "
            "check frontend TOS→records range construction.",
            len(req.records), req.max_items,
        )
    try:
        resp = await run_in_threadpool(
            lambda: generate_quiz_for_topics(
                req.records, max_items=req.max_items, test_labels=req.test_labels
            )
        )
        return resp
    except Exception as e:
        logger.exception("generate error: %s", e)
        raise HTTPException(status_code=500, detail="Generation failed")


@app.post("/generate_from_records")
def generate_from_records_endpoint(payload: dict):
    records = payload.get("records", [])
    if not records:
        return {"items": [], "total_items": 0}
    # Fix-78
    if payload.get("max_items") and len(records) < payload["max_items"]:
        logger.warning(
            "Fix-78 /generate_from_records: got %d records but max_items=%d",
            len(records), payload["max_items"],
        )
    items = generate_questions_from_records(records, max_items=len(records))
    return {"items": items, "total_items": len(items)}


if __name__ == "__main__":
    import argparse
    import uvicorn

    parser = argparse.ArgumentParser()
    parser.add_argument("--jsonl",  "-j", required=False)
    parser.add_argument("--sample", "-n", type=int, default=None)
    parser.add_argument("--serve",  action="store_true")
    args = parser.parse_args()

    if args.jsonl and not args.serve:
        recs = load_jsonl(args.jsonl)
        n = args.sample or min(5, len(recs))
        for r in generate_questions_from_records(recs[:n], max_items=n):
            print(json.dumps(r, indent=2, ensure_ascii=False))
    elif args.serve:
        uvicorn.run("ai_model:app", host="0.0.0.0", port=8000, log_level="info")